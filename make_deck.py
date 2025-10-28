#!/usr/bin/env python3
"""
make_deck.py — Generate a 5-minute PI demo PPTX from a Jira-exported CSV.

Usage:
    python make_deck.py --csv path/to/jira_export.csv \
                        --pi "PI 6" \
                        --top 8 \
                        --out DemoPilot_PI6.pptx \
                        [--pi-number 28] \
                        [--template brand_template.potx] \
                        [--title "DemoPilot — PI 6 Demo"] \
                        [--include-appendix]

What it does:
- Reads a standard Jira CSV export (columns like Issue key, Summary, Issue Type, Status, Priority, Story Points, Assignee, Description, etc.)
- (Optional) Filters to one PI (e.g., "PI 6")
- Computes KPIs (SP totals, assignee contribution)
- Ranks items with a transparent heuristic (no ML): Story Points, “Has Metric Number”, comments/PRs/attachments counters if present, priority, type
- Builds a concise PPTX:
        1) Title slide
        2) KPI slide (SP totals + assignee split)
        3) Top-N item overview
        4) One slide per item (summary, why it matters, metrics cue if found)
        5) (Optional) Appendix table with all items

Dependencies:
    pip install pandas python-pptx
"""

import argparse
import math
import re
import textwrap
from pathlib import Path
from typing import Dict, List, Optional

import pandas as pd
import numpy as np
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.preprocessing import StandardScaler
from sklearn.ensemble import RandomForestRegressor
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR as MSO_ANCHOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE as SHAPE
from pptx.dml.color import RGBColor

# --------------------------
# AI-Enhanced Visual Design System
# --------------------------

class DesignTheme:
    """AI-enhanced design theme with color psychology and typography"""
    
    # Color palette based on data visualization best practices
    PRIMARY_COLORS = {
        'tech_blue': RGBColor(0, 102, 204),      # Trust, reliability
        'success_green': RGBColor(34, 139, 34),   # Success, completion
        'warning_orange': RGBColor(255, 140, 0),  # Attention, priority
        'error_red': RGBColor(220, 20, 60),       # Issues, high priority
        'neutral_gray': RGBColor(64, 64, 64),     # Text, balance
        'accent_purple': RGBColor(138, 43, 226),  # Innovation, creativity
    }
    
    GRADIENT_COLORS = {
        'light_blue': RGBColor(173, 216, 230),
        'medium_blue': RGBColor(100, 149, 237),
        'dark_blue': RGBColor(25, 25, 112),
    }
    
    # Typography hierarchy
    FONT_SIZES = {
        'title': Pt(32),
        'subtitle': Pt(24),
        'heading': Pt(20),
        'body': Pt(16),
        'caption': Pt(12),
        'small': Pt(10)
    }
    
    @staticmethod
    def get_priority_color(priority: str) -> RGBColor:
        """Return color based on priority using color psychology"""
        priority_map = {
            'Highest': DesignTheme.PRIMARY_COLORS['error_red'],
            'High': DesignTheme.PRIMARY_COLORS['warning_orange'],
            'Medium': DesignTheme.PRIMARY_COLORS['tech_blue'],
            'Low': DesignTheme.PRIMARY_COLORS['success_green'],
            'Lowest': DesignTheme.PRIMARY_COLORS['neutral_gray']
        }
        return priority_map.get(priority, DesignTheme.PRIMARY_COLORS['tech_blue'])
    
    @staticmethod
    def get_type_color(issue_type: str) -> RGBColor:
        """Return color based on issue type"""
        type_map = {
            'Story': DesignTheme.PRIMARY_COLORS['tech_blue'],
            'Bug': DesignTheme.PRIMARY_COLORS['error_red'],
            'Task': DesignTheme.PRIMARY_COLORS['success_green'],
            'Spike': DesignTheme.PRIMARY_COLORS['accent_purple'],
            'Epic': DesignTheme.PRIMARY_COLORS['neutral_gray']
        }
        return type_map.get(issue_type, DesignTheme.PRIMARY_COLORS['tech_blue'])

def apply_modern_styling(text_frame, font_size=Pt(16), color=None):
    """Apply modern typography and spacing"""
    for paragraph in text_frame.paragraphs:
        paragraph.font.size = font_size
        paragraph.font.name = 'Segoe UI'
        if color:
            paragraph.font.color.rgb = color
        paragraph.space_after = Pt(6)

def add_visual_separator(slide, y_position=Inches(2.5)):
    """Add a subtle visual separator line"""
    line = slide.shapes.add_connector(
        1, Inches(0.5), y_position, Inches(9.5), y_position
    )
    line.line.color.rgb = DesignTheme.GRADIENT_COLORS['light_blue']
    line.line.width = Pt(1)

def fit_text_frame(tf, max_pt=24, min_pt=10, font_name="Segoe UI", bold=False, center=True):
    """Robust text autoshrink with graceful fallbacks."""
    if tf is None:
        return
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except AttributeError:
        pass
    if center:
        try:
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        except Exception:
            pass

    try:
        tf.fit_text(font=font_name, max_size=max_pt, bold=bold)
        for paragraph in tf.paragraphs:
            if paragraph.font.size and paragraph.font.size.pt > max_pt:
                paragraph.font.size = Pt(max_pt)
        return
    except Exception:
        pass

    lo, hi = min_pt, max_pt
    best = min_pt
    while lo <= hi:
        mid = (lo + hi) // 2
        for paragraph in tf.paragraphs:
            paragraph.font.name = font_name
            paragraph.font.bold = bold or bool(paragraph.font.bold)
            paragraph.font.size = Pt(mid)
            paragraph.word_wrap = True
        total_chars = sum(len(paragraph.text or "") for paragraph in tf.paragraphs)
        needs_shrink = total_chars > 260 or len(tf.paragraphs) > 6
        if needs_shrink:
            hi = mid - 1
        else:
            best = mid
            lo = mid + 1

    for paragraph in tf.paragraphs:
        paragraph.font.size = Pt(best)
        paragraph.font.name = font_name
        paragraph.word_wrap = True


def ensure_slide_text_fits(slide, max_pt=28, min_pt=10):
    """Walk all text frames on a slide and shrink them if needed."""
    for shape in getattr(slide, "shapes", []):
        if not getattr(shape, "has_text_frame", False):
            continue
        tf = shape.text_frame
        if tf is None:
            continue
        cap = max_pt
        for paragraph in tf.paragraphs:
            size = getattr(paragraph.font, "size", None)
            if size:
                try:
                    cap = min(cap, int(size.pt))
                except Exception:
                    pass
        fit_text_frame(tf, max_pt=max(12, cap), min_pt=min_pt)


def create_info_box(slide, left, top, width, height, text, bg_color, text_color=None):
    """Rounded rectangle with autoshrink text (single paragraph)."""
    shape = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = bg_color

    tf = shape.text_frame
    tf.clear()
    paragraph = tf.paragraphs[0]
    paragraph.text = text
    paragraph.font.size = DesignTheme.FONT_SIZES['body']
    paragraph.font.name = 'Segoe UI'
    paragraph.font.bold = True
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.color.rgb = text_color or RGBColor(255, 255, 255)

    tf.margin_left = tf.margin_right = Inches(0.08)
    tf.margin_top = tf.margin_bottom = Inches(0.04)

    fit_text_frame(tf, max_pt=18, min_pt=10, bold=True)
    return shape


def create_metric_card(slide, left, top, width, height, value_text, label_text, bg_color):
    """Two-line KPI card with grouped autoshrink."""
    card = slide.shapes.add_shape(SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = bg_color

    tf = card.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = Inches(0.12)
    tf.margin_top = tf.margin_bottom = Inches(0.06)

    value_paragraph = tf.paragraphs[0]
    value_paragraph.text = value_text
    value_paragraph.font.name = 'Segoe UI'
    value_paragraph.font.bold = True
    value_paragraph.font.size = Pt(30)
    value_paragraph.alignment = PP_ALIGN.CENTER
    value_paragraph.font.color.rgb = RGBColor(255, 255, 255)

    label_paragraph = tf.add_paragraph()
    label_paragraph.text = label_text
    label_paragraph.font.name = 'Segoe UI'
    label_paragraph.font.size = Pt(16)
    label_paragraph.alignment = PP_ALIGN.CENTER
    label_paragraph.font.color.rgb = RGBColor(255, 255, 255)

    fit_text_frame(tf, max_pt=30, min_pt=12, bold=True, center=True)
    return card

# --------------------------
# Helpers: safe field access
# --------------------------

def coerce_num(x, default=0):
    try:
        if pd.isna(x):
            return default
        if isinstance(x, str) and x.strip() == "":
            return default
        return int(float(x))
    except Exception:
        return default

def get(df_row, col, default=""):
    return df_row[col] if col in df_row and pd.notna(df_row[col]) else default

def col_present(df, name):
    return name in df.columns

# --------------------------
# ML-Enhanced Demo Worthiness Scoring
# --------------------------

PRIORITY_W = {"Highest": 2.0, "High": 1.2, "Medium": 0.6, "Low": 0.2, "Lowest": 0.1}
TYPE_W = {"Story": 1.0, "Bug": 0.8, "Task": 0.6, "Spike": 0.4}

# Demo-worthy keywords and patterns
DEMO_KEYWORDS = [
    'user', 'customer', 'interface', 'ui', 'ux', 'dashboard', 'report', 'chart', 'graph',
    'visual', 'display', 'screen', 'page', 'feature', 'functionality', 'improvement',
    'performance', 'optimization', 'automation', 'integration', 'api', 'endpoint',
    'metric', 'analytics', 'data', 'export', 'import', 'notification', 'alert',
    'workflow', 'process', 'enhancement', 'new', 'added', 'implemented', 'delivered'
]

def detect_has_metric(text: str) -> int:
    if not isinstance(text, str):
        return 0
    # detect percentages, deltas, or time mins/hours (e.g., "−30 mins", "≤±5%", "≥4.6/5")
    return 1 if re.search(r"(\d+(\.\d+)?\s*(mins?|hours?)|[±+−-]?\d+%|≥\s*\d+(\.\d+)?/\d+)", text) else 0

def extract_demo_features(row: pd.Series) -> Dict:
    """Extract features that indicate demo worthiness using ML techniques"""
    summary = get(row, "Summary", "").lower()
    desc = get(row, "Description", "").lower()
    combined_text = f"{summary} {desc}"
    
    # Text analysis features
    demo_keyword_count = sum(1 for keyword in DEMO_KEYWORDS if keyword in combined_text)
    
    # Visual/UI related indicators
    ui_indicators = ['ui', 'interface', 'screen', 'page', 'dashboard', 'visual', 'display']
    ui_score = sum(1 for ui in ui_indicators if ui in combined_text)
    
    # User-facing features
    user_facing = ['user', 'customer', 'client', 'end-user', 'stakeholder']
    user_score = sum(1 for user in user_facing if user in combined_text)
    
    # Performance/improvement indicators
    improvement_terms = ['improve', 'enhance', 'optimize', 'faster', 'better', 'new', 'added']
    improvement_score = sum(1 for term in improvement_terms if term in combined_text)
    
    # Measurable impact indicators
    impact_terms = ['reduce', 'increase', 'save', 'time', 'cost', 'efficiency', 'productivity']
    impact_score = sum(1 for term in impact_terms if term in combined_text)
    
    return {
        'demo_keywords': demo_keyword_count,
        'ui_score': ui_score,
        'user_score': user_score,
        'improvement_score': improvement_score,
        'impact_score': impact_score,
        'text_length': len(combined_text),
        'has_screenshot_keywords': 1 if any(term in combined_text for term in ['screenshot', 'image', 'visual', 'demo']) else 0
    }

def score_demo_worthiness(row: pd.Series) -> float:
    """Enhanced scoring using ML features for demo worthiness"""
    # Basic features
    sp = coerce_num(get(row, "Story Points", 0))
    priority = get(row, "Priority", "Medium")
    itype = get(row, "Issue Type", "Story")
    status = get(row, "Status", "")
    
    # For Todo apps and development projects, include "To Do" and "In Progress" items
    # Only exclude cancelled/rejected items
    excluded_statuses = ['cancelled', 'rejected', 'wont do', 'invalid']
    if status.lower() in excluded_statuses:
        return 0.0
    
    # Extract ML features
    demo_features = extract_demo_features(row)
    
    # Enhanced scoring algorithm
    score = 0.0
    
    # Story points weight (higher SP = more significant)
    score += sp * 1.5
    
    # Demo-specific feature weights
    score += demo_features['demo_keywords'] * 2.0
    score += demo_features['ui_score'] * 3.0  # UI changes are highly demo-worthy
    score += demo_features['user_score'] * 2.5  # User-facing features are important
    score += demo_features['improvement_score'] * 1.5
    score += demo_features['impact_score'] * 2.0
    score += demo_features['has_screenshot_keywords'] * 1.0
    
    # Traditional weights
    has_metric = detect_has_metric(get(row, "Description", ""))
    score += has_metric * 2.0
    score += PRIORITY_W.get(priority, 0.5)
    
    # Issue type adjustments for demo worthiness
    demo_type_weights = {
        "Story": 1.2,  # Stories are usually more demo-worthy
        "Bug": 0.6,    # Bugs less so unless high impact
        "Task": 0.4,   # Tasks usually not demo-worthy
        "Spike": 0.2,  # Spikes rarely demo-worthy
        "Epic": 0.1    # Epics too high-level
    }
    score += demo_type_weights.get(itype, 0.6)
    
    # Boost for user-facing stories
    if itype == "Story" and demo_features['user_score'] > 0:
        score *= 1.3
    
    return round(score, 3)

def score_issue(row: pd.Series) -> float:
    """Wrapper for backwards compatibility"""
    return score_demo_worthiness(row)

# --------------------------
# PPT helpers
# --------------------------

def ensure_layout(prs: Presentation, name_contains: str, fallback_index: int = 0):
    for i, layout in enumerate(prs.slide_layouts):
        if name_contains.lower() in layout.name.lower():
            return layout
    return prs.slide_layouts[fallback_index]

def add_title_slide(prs: Presentation, title_text: str, subtitle_text: str):
    """Create an AI-enhanced title slide with modern design"""
    layout = ensure_layout(prs, "blank", 6)  # Use blank layout for full control
    slide = prs.slides.add_slide(layout)
    
    # Background gradient effect (simulated with shapes)
    bg_shape = slide.shapes.add_shape(
        SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(7.5)
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = DesignTheme.GRADIENT_COLORS['light_blue']
    bg_shape.line.fill.background()
    
    # Accent stripe
    accent_shape = slide.shapes.add_shape(
        SHAPE.RECTANGLE, Inches(0), Inches(6.5), Inches(10), Inches(1)
    )
    accent_shape.fill.solid()
    accent_shape.fill.fore_color.rgb = DesignTheme.PRIMARY_COLORS['tech_blue']
    accent_shape.line.fill.background()
    
    # Main title with modern typography
    title_box = slide.shapes.add_textbox(
        Inches(1), Inches(2), Inches(8), Inches(2)
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = title_text
    title_p.font.size = DesignTheme.FONT_SIZES['title']
    title_p.font.name = 'Segoe UI'
    title_p.font.bold = True
    title_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
    title_p.alignment = PP_ALIGN.CENTER
    fit_text_frame(title_frame, max_pt=32, bold=True)
    
    # Subtitle with accent color
    subtitle_box = slide.shapes.add_textbox(
        Inches(1), Inches(4.2), Inches(8), Inches(1)
    )
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.clear()
    subtitle_p = subtitle_frame.paragraphs[0]
    subtitle_p.text = subtitle_text
    subtitle_p.font.size = DesignTheme.FONT_SIZES['subtitle']
    subtitle_p.font.name = 'Segoe UI'
    subtitle_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['tech_blue']
    subtitle_p.alignment = PP_ALIGN.CENTER
    fit_text_frame(subtitle_frame, max_pt=22)

    ensure_slide_text_fits(slide, max_pt=28, min_pt=10)

def add_kpi_slide(prs: Presentation, pi_name: str, df: pd.DataFrame):
    """Create an AI-enhanced KPI slide with visual elements"""
    layout = ensure_layout(prs, "blank", 6)
    slide = prs.slides.add_slide(layout)
    
    # Modern title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = f"{pi_name} — PI Analytics & KPIs"
    title_p.font.size = DesignTheme.FONT_SIZES['title']
    title_p.font.name = 'Segoe UI'
    title_p.font.bold = True
    title_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
    
    # Compute metrics with flexible column handling
    total_sp = int(df["Story Points"].fillna(0).apply(coerce_num).sum())
    by_type = df.groupby("Issue Type")["Story Points"].sum().to_dict()
    
    # Handle optional Assignee column
    by_assignee = {}
    if "Assignee" in df.columns:
        by_assignee = df.groupby("Assignee")["Story Points"].sum().sort_values(ascending=False).to_dict()
    elif "Components" in df.columns:
        by_assignee = df.groupby("Components")["Story Points"].sum().sort_values(ascending=False).to_dict()
    
    total_items = len(df)
    
    # Key metrics cards in top row (no overflow)
    metrics = [
        ("Total Story Points", str(total_sp), DesignTheme.PRIMARY_COLORS['tech_blue']),
        ("Items Completed", str(total_items), DesignTheme.PRIMARY_COLORS['success_green']),
        ("Avg SP/Item", f"{(total_sp/max(total_items,1)):.1f}", DesignTheme.PRIMARY_COLORS['accent_purple'])
    ]
    card_width = Inches(2.8)
    card_height = Inches(1.2)
    start_x = Inches(0.7)

    for i, (label, value, color) in enumerate(metrics):
        x_pos = start_x + i * (card_width + Inches(0.3))
        create_metric_card(
            slide,
            x_pos, Inches(1.3),
            card_width, card_height,
            value_text=value,
            label_text=label,
            bg_color=color
        )
    
    # Issue type breakdown with color coding
    if by_type:
        type_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3), Inches(4), Inches(3)
        )
        type_frame = type_box.text_frame
        type_frame.clear()
        
        # Section header
        header_p = type_frame.paragraphs[0]
        header_p.text = "📊 Breakdown by Issue Type"
        header_p.font.size = DesignTheme.FONT_SIZES['heading']
        header_p.font.name = 'Segoe UI'
        header_p.font.bold = True
        header_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
        
        # Type items with visual indicators
        for issue_type, sp_count in by_type.items():
            p = type_frame.add_paragraph()
            percentage = (sp_count / total_sp) * 100 if total_sp > 0 else 0
            p.text = f"▎ {issue_type}: {int(sp_count)} SP ({percentage:.1f}%)"
            p.font.size = DesignTheme.FONT_SIZES['body']
            p.font.name = 'Segoe UI'
            p.font.color.rgb = DesignTheme.get_type_color(issue_type)
            p.level = 1
        fit_text_frame(type_frame, max_pt=18)
    
    # Team contribution with enhanced visuals
    if by_assignee:
        team_box = slide.shapes.add_textbox(
            Inches(5), Inches(3), Inches(4), Inches(3)
        )
        team_frame = team_box.text_frame
        team_frame.clear()
        
        # Section header - adapt based on data type
        header_text = "👥 Team Contribution" if "Assignee" in df.columns else "🏗️ Component Breakdown"
        header_p = team_frame.paragraphs[0]
        header_p.text = header_text
        header_p.font.size = DesignTheme.FONT_SIZES['heading']
        header_p.font.name = 'Segoe UI'
        header_p.font.bold = True
        header_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
        
        # Team member contributions with progress bars simulation
        colors = [DesignTheme.PRIMARY_COLORS['tech_blue'], 
                 DesignTheme.PRIMARY_COLORS['success_green'],
                 DesignTheme.PRIMARY_COLORS['warning_orange'],
                 DesignTheme.PRIMARY_COLORS['accent_purple']]
        
        for i, (assignee, sp_count) in enumerate(list(by_assignee.items())[:4]):
            p = team_frame.add_paragraph()
            percentage = (sp_count / total_sp) * 100 if total_sp > 0 else 0
            # Create visual progress indicator
            bar_length = int(percentage / 10)  # Scale to 10 chars max
            progress_bar = "█" * bar_length + "░" * (10 - bar_length)
            p.text = f"{assignee}: {int(sp_count)} SP\n{progress_bar} {percentage:.1f}%"
            p.font.size = DesignTheme.FONT_SIZES['body']
            p.font.name = 'Segoe UI'
            p.font.color.rgb = colors[i % len(colors)]
            p.level = 1
        fit_text_frame(team_frame, max_pt=18)
    
    # Add visual separator
    add_visual_separator(slide, Inches(6.2))

    ensure_slide_text_fits(slide, max_pt=28, min_pt=10)

def add_overview_slide(
    prs: Presentation,
    pi_name: str,
    top_items: List[Dict],
    page_number: int = 1,
    total_pages: int = 1,
    start_rank: int = 1,
    total_items: Optional[int] = None,
):
    """Create an AI-enhanced overview slide with visual hierarchy.

    When total_pages > 1 we suffix the title with a page indicator and continue
    numbering using the provided start_rank so the roadmap order stays clear."""
    layout = ensure_layout(prs, "blank", 6)
    slide = prs.slides.add_slide(layout)
    
    # Modern title with subtitle
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(1)
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    slide_title = f"{pi_name} — PI Demo Roadmap"
    if total_pages > 1:
        slide_title = f"{slide_title} — Page {page_number}"
    title_p.text = slide_title
    title_p.font.size = DesignTheme.FONT_SIZES['title']
    title_p.font.name = 'Segoe UI'
    title_p.font.bold = True
    title_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
    
    # Subtitle
    subtitle_p = title_frame.add_paragraph()
    overall_items = total_items or len(top_items)
    if total_pages > 1:
        end_rank = start_rank + len(top_items) - 1
        subtitle_p.text = f"Items {start_rank}-{end_rank} of {overall_items}"
    else:
        subtitle_p.text = f"Top {overall_items} Demo-Worthy Items"
    subtitle_p.font.size = DesignTheme.FONT_SIZES['subtitle']
    subtitle_p.font.name = 'Segoe UI'
    subtitle_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['tech_blue']
    
    # Create visual item cards
    start_y = Inches(1.75)
    card_spacing = Inches(0.16)

    num_items = len(top_items)
    columns = 2 if num_items > 6 else 1
    items_per_col = math.ceil(num_items / columns)
    column_gap = Inches(0.3)
    column_width = Inches(4.35) if columns == 2 else Inches(8.5)

    # Dynamically size cards to fit available vertical space
    max_body_height = Inches(6.1) - start_y
    min_card_height = Inches(0.5)
    max_card_height = Inches(0.88)
    if items_per_col > 0:
        available_height = max_body_height - (items_per_col - 1) * card_spacing
        card_height = max(min_card_height, min(max_card_height, available_height / items_per_col))
    else:
        card_height = min_card_height

    for i, item in enumerate(top_items):
        col_idx = i // items_per_col
        row_idx = i % items_per_col
        x_pos = Inches(0.5) + col_idx * (column_width + column_gap)
        y_pos = start_y + row_idx * (card_height + card_spacing)

        # Priority color coding
        priority_color = DesignTheme.get_priority_color(item.get('Priority', 'Medium'))
        type_color = DesignTheme.get_type_color(item.get('Issue Type', 'Story'))

        # Main item card
        item_box = slide.shapes.add_textbox(
            x_pos, y_pos, column_width, card_height
        )
        item_frame = item_box.text_frame
        item_frame.clear()
        item_frame.margin_left = Inches(0.05)
        item_frame.margin_right = Inches(0.05)
        item_frame.margin_top = Inches(0.04)
        item_frame.margin_bottom = Inches(0.02)

        # Item header with number and key — allow wrapping and autoshrink per card
        demo_score = item.get('demo_score', 0)
        header_p = item_frame.paragraphs[0]
        # keep full-ish summary but cap to a generous length; rely on autoshrink to fit
        header_p.text = f"{start_rank + i}. {item['Issue key']} — {_shorten(item.get('Summary',''), 140)}"
        header_p.font.size = Pt(16)
        header_p.font.name = 'Segoe UI'
        header_p.font.bold = True
        header_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
        # ensure word wrap and allow autoshrink per-card
        header_p.word_wrap = True

        # Metadata line with visual indicators
        meta_p = item_frame.add_paragraph()
        sp = item.get('Story Points', 0)
        priority = item.get('Priority', 'Medium')
        issue_type = item.get('Issue Type', 'Story')
        meta_text = f"   🏷️ {issue_type} | ⚡ {priority} | 📊 {sp} SP"
        if demo_score:
            meta_text += f" | 🎯 Score {demo_score:.1f}"
        meta_p.text = meta_text
        meta_p.font.size = Pt(13)
        meta_p.font.name = 'Segoe UI'
        meta_p.font.color.rgb = type_color
        meta_p.level = 1

        # Priority indicator strip
        priority_strip = slide.shapes.add_shape(
            SHAPE.RECTANGLE, x_pos + column_width + Inches(0.05), y_pos, Inches(0.12), card_height
        )
        priority_strip.fill.solid()
        priority_strip.fill.fore_color.rgb = priority_color
        priority_strip.line.fill.background()

        # autoshrink text inside this specific card to avoid squishing elsewhere
        try:
            fit_text_frame(item_frame, max_pt=16, min_pt=10)
        except Exception:
            pass
    
    # Add legend at bottom
    legend_y = start_y + items_per_col * (card_height + card_spacing) + Inches(0.25)
    legend_y = min(legend_y, Inches(6.6))
    legend_box = slide.shapes.add_textbox(
        Inches(0.5), legend_y, Inches(9), Inches(0.45)
    )
    legend_frame = legend_box.text_frame
    legend_frame.clear()
    legend_p = legend_frame.paragraphs[0]
    legend_p.text = "💡 Highlights chosen for user impact, visuals, and business value"
    legend_p.font.size = DesignTheme.FONT_SIZES['caption']
    legend_p.font.name = 'Segoe UI'
    legend_p.font.italic = True
    legend_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
    legend_p.alignment = PP_ALIGN.CENTER
    fit_text_frame(legend_frame, max_pt=16)

    ensure_slide_text_fits(slide, max_pt=28, min_pt=10)

def _shorten(text: str, limit=220):
    if not text:
        return ""
    t = re.sub(r"\s+", " ", text).strip()
    if len(t) <= limit:
        return t
    return t[:limit - 1] + "…"

def _bullets_from_description(desc: str) -> List[str]:
    if not isinstance(desc, str) or not desc.strip():
        return []
    # try to pull “Why/Value/Demo/AC” lines if present
    bullets = []
    for label in ["Why", "Value", "Metric", "Demo", "Acceptance", "Risk", "Context", "Note"]:
        m = re.search(rf"(?i)^{label}[:\s-]+(.+)$", desc, flags=re.MULTILINE)
        if m:
            bullets.append(f"{label}: {m.group(1).strip()}")
    # fallbacks: first two sentences
    if not bullets:
        sents = re.split(r"(?<=[.!?])\s+", desc.strip())
        bullets = [s for s in sents[:3]]
    # keep them short
    return [_shorten(b, 120) for b in bullets[:4]]


def normalize_pi_label(value: str) -> str:
    label = "" if value is None else str(value).strip()
    if not label:
        return "PI"
    if label.lower().startswith("pi"):
        remainder = label[2:].strip()
        return f"PI {remainder}" if remainder else "PI"
    match = re.search(r"(\d+)", label)
    if match:
        return f"PI {match.group(1)}"
    return f"PI {label}"

def add_item_slide(prs: Presentation, item: Dict, idx: int, total: int):
    """Create an AI-enhanced individual item slide with modern design"""
    layout = ensure_layout(prs, "blank", 6)
    slide = prs.slides.add_slide(layout)

    # Header section with progress indicator
    header_height = Inches(1.2)
    
    # Progress bar background
    progress_bg = slide.shapes.add_shape(
        SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), header_height
    )
    progress_bg.fill.solid()
    progress_bg.fill.fore_color.rgb = DesignTheme.GRADIENT_COLORS['light_blue']
    progress_bg.line.fill.background()
    
    # Progress indicator
    progress_width = (idx / total) * Inches(10)
    progress_bar = slide.shapes.add_shape(
        SHAPE.RECTANGLE, Inches(0), Inches(0), progress_width, Inches(0.1)
    )
    progress_bar.fill.solid()
    progress_bar.fill.fore_color.rgb = DesignTheme.PRIMARY_COLORS['tech_blue']
    progress_bar.line.fill.background()
    
    # Title with modern formatting
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.2), Inches(8), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = f"{item['Issue key']}: {_shorten(item['Summary'], 60)}"
    title_p.font.size = DesignTheme.FONT_SIZES['title']
    title_p.font.name = 'Segoe UI'
    title_p.font.bold = True
    title_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
    
    # Progress indicator text
    progress_text = slide.shapes.add_textbox(
        Inches(8.5), Inches(0.2), Inches(1.3), Inches(0.8)
    )
    progress_frame = progress_text.text_frame
    progress_frame.clear()
    progress_p = progress_frame.paragraphs[0]
    progress_p.text = f"{idx}/{total}"
    progress_p.font.size = DesignTheme.FONT_SIZES['heading']
    progress_p.font.name = 'Segoe UI'
    progress_p.font.bold = True
    progress_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['tech_blue']
    progress_p.alignment = PP_ALIGN.RIGHT
    fit_text_frame(progress_frame, max_pt=20, bold=True)
    
    # Metadata cards section
    card_y = Inches(1.5)
    card_width = Inches(1.8)
    card_height = Inches(0.6)
    
    # Create metadata cards with flexible data
    metadata_cards = [
        (item.get('Issue Type', 'Story'), DesignTheme.get_type_color(item.get('Issue Type', 'Story'))),
        (item.get('Priority', 'Medium'), DesignTheme.get_priority_color(item.get('Priority', 'Medium'))),
        (f"{item.get('Story Points', 0)} SP", DesignTheme.PRIMARY_COLORS['accent_purple']),
        (item.get('Status', 'To Do'), DesignTheme.PRIMARY_COLORS['success_green'])
    ]
    
    for i, (text, color) in enumerate(metadata_cards):
        x_pos = Inches(0.5) + i * (card_width + Inches(0.3))
        create_info_box(slide, x_pos, card_y, card_width, card_height, text, color)
    
    # Demo worthiness explanation with enhanced styling
    demo_reasons = _get_demo_reasons(item)
    if demo_reasons:
        demo_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.4), Inches(9), Inches(0.8)
        )
        demo_frame = demo_box.text_frame
        demo_frame.clear()
        demo_p = demo_frame.paragraphs[0]
        demo_p.text = f"🎯 Why it matters: {demo_reasons}"
        demo_p.font.size = DesignTheme.FONT_SIZES['body']
        demo_p.font.name = 'Segoe UI'
        demo_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['success_green']
        demo_p.font.italic = True
        fit_text_frame(demo_frame, max_pt=18)
    
    # Content section with better formatting
    content_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(3.5), Inches(9), Inches(3.5)
    )
    content_frame = content_box.text_frame
    content_frame.clear()
    
    # Section header
    header_p = content_frame.paragraphs[0]
    header_p.text = "📋 Key Details & Demo Points"
    header_p.font.size = DesignTheme.FONT_SIZES['heading']
    header_p.font.name = 'Segoe UI'
    header_p.font.bold = True
    header_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
    
    # Enhanced bullets from description
    bullets = _bullets_from_description(item.get("Description", ""))
    if not bullets:
        default_bullets = []
        if item.get('Assignee'):
            default_bullets.append(f"Assignee: {item['Assignee']}")
        if item.get('components'):
            default_bullets.append(f"Components: {item['components']}")
        elif item.get('Components'):
            default_bullets.append(f"Components: {item['Components']}")
        default_bullets.append("Ready for demonstration")
        bullets = default_bullets
    
    for bullet in bullets:
        trimmed_bullet = _shorten(bullet, 120)
        bp = content_frame.add_paragraph()
        bp.text = f"▸ {trimmed_bullet}"
        bp.font.size = DesignTheme.FONT_SIZES['body']
        bp.font.name = 'Segoe UI'
        bp.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
        bp.level = 1
        bp.space_after = Pt(8)
    fit_text_frame(content_frame, max_pt=18)
    
    # Add visual separator
    add_visual_separator(slide, Inches(7.2))

    ensure_slide_text_fits(slide, max_pt=28, min_pt=10)

def _get_demo_reasons(item: Dict) -> str:
    """Generate explanation for why this item is demo-worthy"""
    reasons = []
    summary = item.get("Summary", "").lower()
    desc = item.get("Description", "").lower()
    combined = f"{summary} {desc}"
    
    # Check for UI/visual elements
    if any(term in combined for term in ['ui', 'interface', 'dashboard', 'chart', 'visual', 'screen']):
        reasons.append("visual/UI changes")
    
    # Check for user-facing features
    if any(term in combined for term in ['user', 'customer', 'feature', 'functionality']):
        reasons.append("user-facing feature")
    
    # Check for performance/improvements
    if any(term in combined for term in ['improve', 'optimize', 'performance', 'faster', 'better']):
        reasons.append("performance improvement")
    
    # Check for high story points
    sp = item.get("Story Points", 0)
    if isinstance(sp, (int, float)) and sp >= 5:
        reasons.append(f"high impact ({sp} SP)")
    
    # Check for priority
    priority = item.get("Priority", "")
    if priority in ["High", "Highest"]:
        reasons.append("high priority")
    
    return ", ".join(reasons[:3]) if reasons else "completed deliverable"

def add_appendix_table(prs: Presentation, df: pd.DataFrame, title="Appendix — Complete PI Overview"):
    """Create an AI-enhanced appendix table with visual styling"""
    layout = ensure_layout(prs, "blank", 6)
    slide = prs.slides.add_slide(layout)
    
    # Modern title
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.3), Inches(9), Inches(0.8)
    )
    title_frame = title_box.text_frame
    title_frame.clear()
    title_p = title_frame.paragraphs[0]
    title_p.text = title
    title_p.font.size = DesignTheme.FONT_SIZES['title']
    title_p.font.name = 'Segoe UI'
    title_p.font.bold = True
    title_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
    fit_text_frame(title_frame, max_pt=30, bold=True)
    
    # Summary stats
    stats_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1), Inches(9), Inches(0.4)
    )
    stats_frame = stats_box.text_frame
    stats_frame.clear()
    stats_p = stats_frame.paragraphs[0]
    total_items = len(df)
    total_sp = int(df["Story Points"].fillna(0).apply(coerce_num).sum())
    stats_p.text = f"📊 Total: {total_items} items | {total_sp} story points"
    stats_p.font.size = DesignTheme.FONT_SIZES['body']
    stats_p.font.name = 'Segoe UI'
    stats_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['tech_blue']
    stats_p.alignment = PP_ALIGN.CENTER
    fit_text_frame(stats_frame, max_pt=18)
    
    # Enhanced table with better columns - flexible for different CSV structures
    cols_pref = ["Issue key", "Summary", "Issue Type", "Priority", "Story Points", "Status"]
    # Add Assignee or Components if available
    if "Assignee" in df.columns:
        cols_pref.insert(-2, "Assignee")
    elif "Components" in df.columns:
        cols_pref.insert(-2, "Components")
    
    cols = [c for c in cols_pref if c in df.columns]
    if not cols:
        return
    
    # Limit rows for readability
    display_df = df.head(12)
    num_rows = len(display_df) + 1  # header + data rows
    num_cols = len(cols)

    x, y, cx, cy = Inches(0.3), Inches(1.8), Inches(9.4), Inches(5)
    table_shape = slide.shapes.add_table(num_rows, num_cols, x, y, cx, cy)
    table = table_shape.table
    
    # Style header row
    for j, col in enumerate(cols):
        cell = table.cell(0, j)
        cell.text = col
        # Header styling
        cell.fill.solid()
        cell.fill.fore_color.rgb = DesignTheme.PRIMARY_COLORS['tech_blue']
        
        # Header text styling
        paragraph = cell.text_frame.paragraphs[0]
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
        paragraph.font.bold = True
        paragraph.font.size = DesignTheme.FONT_SIZES['body']
        paragraph.font.name = 'Segoe UI'
        paragraph.alignment = PP_ALIGN.CENTER
        fit_text_frame(cell.text_frame, max_pt=16, bold=True)

    # Data rows with alternating colors and smart formatting
    for i in range(1, num_rows):
        if i-1 >= len(display_df):
            break
        row = display_df.iloc[i-1]
        
        # Alternating row colors
        row_color = DesignTheme.GRADIENT_COLORS['light_blue'] if i % 2 == 0 else RGBColor(255, 255, 255)
        
        for j, col in enumerate(cols):
            cell = table.cell(i, j)
            val = str(row[col]) if col in row and pd.notna(row[col]) else ""
            
            # Smart text formatting
            if col == "Summary":
                val = _shorten(val, 50)
            elif col == "Issue key":
                val = val  # Keep full key
            elif col == "Story Points":
                val = str(int(float(val))) if val and val != "0.0" else "0"
            
            cell.text = val
            
            # Cell styling
            cell.fill.solid()
            cell.fill.fore_color.rgb = row_color
            
            # Text styling based on content
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.font.size = DesignTheme.FONT_SIZES['small']
            paragraph.font.name = 'Segoe UI'
            
            # Color coding based on column type
            if col == "Priority":
                paragraph.font.color.rgb = DesignTheme.get_priority_color(val)
                paragraph.font.bold = True
            elif col == "Issue Type":
                paragraph.font.color.rgb = DesignTheme.get_type_color(val)
            elif col == "Status" and val.lower() in ['done', 'closed', 'resolved']:
                paragraph.font.color.rgb = DesignTheme.PRIMARY_COLORS['success_green']
                paragraph.font.bold = True
            else:
                paragraph.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']

            fit_text_frame(cell.text_frame, max_pt=13)
    
    # Add footer note if table is truncated
    if len(df) > 12:
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7), Inches(9), Inches(0.3)
        )
        footer_frame = footer_box.text_frame
        footer_frame.clear()
        footer_p = footer_frame.paragraphs[0]
        footer_p.text = f"Note: Showing first 12 of {len(df)} total items. Complete data available in source CSV."
        footer_p.font.size = DesignTheme.FONT_SIZES['caption']
        footer_p.font.name = 'Segoe UI'
        footer_p.font.italic = True
        footer_p.font.color.rgb = DesignTheme.PRIMARY_COLORS['neutral_gray']
        footer_p.alignment = PP_ALIGN.CENTER
        fit_text_frame(footer_frame, max_pt=14)

    ensure_slide_text_fits(slide, max_pt=28, min_pt=10)

# --------------------------
# Main pipeline
# --------------------------

def rank_dataframe(df: pd.DataFrame, top_n: int) -> List[Dict]:
    """Enhanced ranking for demo-worthy items using ML-based scoring"""
    if "Status" in df.columns:
        dfc = df.copy()
    else:
        dfc = df.copy()
        dfc["Status"] = "Done"
    
    # Add Issue key if not present (use Summary or create index)
    if "Issue key" not in dfc.columns:
        if "Summary" in dfc.columns:
            # Create issue keys from index for Todo app
            dfc["Issue key"] = [f"TODO-{i+1}" for i in range(len(dfc))]
        else:
            dfc["Issue key"] = [f"ITEM-{i+1}" for i in range(len(dfc))]

    # Filter for demo-worthy items (more inclusive for development projects)
    excluded_statuses = ['cancelled', 'rejected', 'wont do', 'invalid']
    if "Status" in dfc.columns:
        status_filter = ~dfc["Status"].str.lower().isin(excluded_statuses)
    else:
        status_filter = True
    
    type_filter = ~dfc["Issue Type"].str.lower().isin(["epic"])  # Include spikes for Todo apps
    
    candidates = dfc[status_filter & type_filter].copy()
    
    if len(candidates) == 0:
        print("Warning: No demo-worthy items found")
        return []

    # Compute demo worthiness score
    candidates["__demo_score__"] = candidates.apply(score_demo_worthiness, axis=1)
    
    # Filter out items with very low demo scores (less likely to be demo-worthy)
    min_demo_score = 1.0  # Minimum threshold for demo worthiness
    demo_worthy = candidates[candidates["__demo_score__"] >= min_demo_score].copy()
    
    if len(demo_worthy) == 0:
        print(f"Warning: No items meet demo worthiness threshold ({min_demo_score})")
        print("Top 5 candidate items:")
        top_scores = candidates.nlargest(5, "__demo_score__")
        for _, row in top_scores.iterrows():
            print(f"  {row['Issue key']}: {row['Summary'][:50]}...")
        # Fall back to top items even if below threshold
        demo_worthy = candidates.nlargest(min(top_n, len(candidates)), "__demo_score__")
    
    # Sort by demo score, then by story points, then by updated date
    sort_cols = ["__demo_score__"]
    ascending = [False]
    if "Story Points" in demo_worthy.columns:
        sort_cols.append("Story Points")
        ascending.append(False)
    if "Updated" in demo_worthy.columns:
        sort_cols.append("Updated")
        ascending.append(False)

    demo_worthy = demo_worthy.sort_values(sort_cols, ascending=ascending)
    
    # Ensure diversity in issue types if possible
    top_items = []
    used_types = set()
    
    # First pass: get diverse types
    for _, row in demo_worthy.iterrows():
        if len(top_items) >= top_n:
            break
        issue_type = row["Issue Type"]
        if issue_type not in used_types or len(used_types) >= 3:
            top_items.append(row.to_dict())
            used_types.add(issue_type)
    
    # Second pass: fill remaining slots with highest scoring items
    for _, row in demo_worthy.iterrows():
        if len(top_items) >= top_n:
            break
        row_dict = row.to_dict()
        if row_dict not in top_items:
            top_items.append(row_dict)
    
    # Pack metadata and clean up
    for r in top_items:
        r["components"] = r.get("Components", "")
        r["demo_score"] = r.get("__demo_score__", 0)
        # Remove internal scoring columns
        r.pop("__demo_score__", None)
    
    print(f"Selected {len(top_items)} demo-worthy items:")
    for i, item in enumerate(top_items, 1):
        print(f"  {i}. {item['Issue key']}: {item['Summary'][:60]}...")
    
    return top_items

def build_deck(
    csv_path: Path,
    out_path: Path,
    pi_filter: str = None,
    pi_number: str = None,
    top_n: int = 8,
    template_path: Path = None,
    title: str = None,
    include_appendix: bool = False
):
    df = pd.read_csv(csv_path, dtype=str).fillna("")

    # normalize numeric story points
    if "Story Points" in df.columns:
        df["Story Points"] = df["Story Points"].apply(coerce_num)
    else:
        df["Story Points"] = 0

    # PI filter
    if "PI" in df.columns:
        pi_column = "PI"
    elif "Sprint" in df.columns:
        pi_column = "Sprint"
    else:
        pi_column = None

    pi_label = pi_filter
    if pi_column:
        if not pi_label:
            col_series = df[pi_column].dropna()
            pi_label = col_series.mode().iat[0] if len(col_series) else "PI"
        if pi_filter:
            df = df[df[pi_column] == pi_filter].copy()
    else:
        pi_label = pi_filter or "PI"

    if pi_number:
        pi_name = normalize_pi_label(pi_number)
    else:
        pi_name = normalize_pi_label(pi_label)

    if title:
        title_clean = str(title).strip()
        if re.match(r"(?i)^pi[\s-]*\d+$", title_clean):
            pi_name = normalize_pi_label(title_clean)

    # create presentation
    prs = Presentation(str(template_path)) if template_path and Path(template_path).exists() else Presentation()

    # Title slide
    main_title = title or f"DemoPilot — {pi_name} Demo"
    subtitle = f"Auto-generated from Jira CSV • Top {top_n} items"
    add_title_slide(prs, main_title, subtitle)

    # KPI slide
    if len(df):
        add_kpi_slide(prs, pi_name, df)

    # Ranking
    top_items = rank_dataframe(df, top_n=top_n)
    if top_items:
        items_per_page = 5
        total_items = len(top_items)
        total_pages = math.ceil(total_items / items_per_page)
        for page_index in range(total_pages):
            start = page_index * items_per_page
            end = start + items_per_page
            page_items = top_items[start:end]
            add_overview_slide(
                prs,
                pi_name,
                page_items,
                page_number=page_index + 1,
                total_pages=total_pages,
                start_rank=start + 1,
                total_items=total_items,
            )
        for i, item in enumerate(top_items, 1):
            add_item_slide(prs, item, i, len(top_items))

    # Appendix (optional)
    if include_appendix:
        add_appendix_table(prs, df)

    prs.save(str(out_path))

# --------------------------
# CLI
# --------------------------

def parse_args():
    ap = argparse.ArgumentParser(description="Generate a PPTX PI demo from a Jira CSV export.")
    ap.add_argument("--csv", required=True, help="Path to Jira CSV export")
    ap.add_argument("--out", required=True, help="Output PPTX path")
    ap.add_argument("--pi", "--sprint", dest="pi", default=None, help="Filter to this PI name (e.g., 'PI 6')")
    ap.add_argument("--pi-number", default=None, help="Override the displayed PI label (e.g., '28' or 'PI 28')")
    ap.add_argument("--top", type=int, default=8, help="Top-N items to include")
    ap.add_argument("--template", default=None, help="Optional .potx/.pptx to use as brand template")
    ap.add_argument("--title", default=None, help="Custom deck title")
    ap.add_argument("--include-appendix", action="store_true", help="Add an appendix slide with an items table")
    return ap.parse_args()

if __name__ == "__main__":
    args = parse_args()
    build_deck(
        csv_path=Path(args.csv),
        out_path=Path(args.out),
    pi_filter=args.pi,
    pi_number=args.pi_number,
        top_n=args.top,
        template_path=Path(args.template) if args.template else None,
        title=args.title,
        include_appendix=args.include_appendix,
    )
